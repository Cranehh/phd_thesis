# -*- coding: utf-8 -*-
"""spExample 中文版：在问卷绘图.pptx 的SP情境(含MaaS)幻灯片上把文字整段译为中文后导出。"""
from pptx import Presentation
p=Presentation('/mnt/user-data/uploads/2023_04MaaS单次出行/问卷绘图.pptx')
D={
 'With the addition of four MaaS options, which one would you prefer for your current trip?':'在新增四种MaaS选项后，您更愿意选择哪一种完成本次出行？',
 'Which mode will you choose to make the trip in the above context?':'在上述情境下，您会选择哪种方式完成本次出行？',
 'Trip purpose: Non-commuting':'出行目的：非通勤','Departure time: 8:00-10:00':'出发时间：8:00–10:00',
 'Taxi/ Ride-sourcing':'出租/网约车','Taxi/Ride-sourcing':'出租/网约车',
 'Metro & Shared bike':'地铁+共享单车','Metro & Taxi/ ride-sourcing':'地铁+出租/网约车','Metro & Bus':'地铁+公交',
 'In-taxi/ Ride-sourcing time':'出租/网约车内时间','In-metro time':'地铁内时间','In-bus time':'公交内时间',
 'In-vehicle time':'车内时间','Waiting time':'等待时间','Walking distance':'步行距离','Riding distance':'骑行距离',
 'Parking and fuel costs':'停车与油耗费用','Travel time':'出行时间','Travel cost':'出行费用',
 'Trip purpose:':'出行目的：','Departure time:':'出发时间：','Non-commuting':'非通勤',
 'MaaS Option 1':'MaaS方案 M1','MaaS Option 2':'MaaS方案 M2','MaaS Option 3':'MaaS方案 M3','MaaS Option 4':'MaaS方案 M4',
 'Ride-sharing':'合乘','Destination':'终点','Origin':'起点','Private car':'私家车','Bus':'公交',
 'min':'分钟','CNY':'元',
}
keys=sorted(D,key=len,reverse=True)
def tr(s):
    for k in keys: s=s.replace(k,D[k])
    return s
for sh in p.slides[0].shapes:
    if not sh.has_text_frame: continue
    for para in sh.text_frame.paragraphs:
        runs=para.runs
        if not runs: continue
        full=''.join(r.text for r in runs)
        if full.strip():
            runs[0].text=tr(full)
            for r in runs[1:]: r.text=''
p.save('ppt/spExample_cn.pptx'); print('saved')
