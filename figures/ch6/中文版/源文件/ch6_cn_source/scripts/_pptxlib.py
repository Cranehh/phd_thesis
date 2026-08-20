# -*- coding: utf-8 -*-
"""python-pptx 构图小工具（中文框架图专用）。

坐标一律用「英寸」，(0,0) 在左上角。字体：拉丁文 Times New Roman，
中文 East-Asian 用 Noto Serif CJK SC（≈宋体）。Windows 上可在 PowerPoint
里全选改成「宋体」。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

LATIN = 'Times New Roman'
EA = 'Noto Serif CJK SC'

# 调色（与英文原图接近的中性配色）
GRAY_FILL = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_PANEL = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARKLINE = RGBColor(0x2E, 0x54, 0x86)   # 深蓝描边（原图框线）
BLUE_BG = RGBColor(0xDC, 0xE6, 0xF2)    # 下半部淡蓝底
BLUE_TAB = RGBColor(0xC6, 0xD9, 0xF0)
ORANGE_TAB = RGBColor(0xFB, 0xE4, 0xD5)


def new_prez(w_in, h_in):
    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    return prs, slide


def _set_run_font(run, size, bold, color, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = LATIN
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', LATIN if tag == 'a:latin' else EA)


def _fill_line(shape, fill, line_color, line_w=1.0, dashed=False):
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
        if dashed:
            ln = shape.line._get_or_add_ln()
            d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
            ln.append(d)
    shape.shadow.inherit = False


def box(slide, x, y, w, h, text='', fill=WHITE, line_color=DARKLINE,
        line_w=1.0, size=11, bold=False, color=BLACK, align='c',
        anchor='m', shape=MSO_SHAPE.ROUNDED_RECTANGLE, dashed=False,
        italic=False, wrap=True):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    _fill_line(sp, fill, line_color, line_w, dashed)
    # rounded-rect 圆角调小一点
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = 0.08
        except Exception:
            pass
    tf = sp.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE,
                          'b': MSO_ANCHOR.BOTTOM}[anchor]
    for m in ('margin_left', 'margin_right'):
        setattr(tf, m, Emu(18000))
    for m in ('margin_top', 'margin_bottom'):
        setattr(tf, m, Emu(9000))
    _write(tf.paragraphs[0], text, size, bold, color, align, italic)
    return sp


def textbox(slide, x, y, w, h, text='', size=11, bold=False, color=BLACK,
            align='c', anchor='m', italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE,
                          'b': MSO_ANCHOR.BOTTOM}[anchor]
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    _write(tf.paragraphs[0], text, size, bold, color, align, italic)
    return tb


def _write(p, text, size, bold, color, align, italic=False):
    p.alignment = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER,
                   'r': PP_ALIGN.RIGHT}[align]
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        if i == 0:
            run = p.add_run()
            run.text = ln
            _set_run_font(run, size, bold, color, italic)
        else:
            p2 = p._parent.add_paragraph()
            p2.alignment = p.alignment
            run = p2.add_run()
            run.text = ln
            _set_run_font(run, size, bold, color, italic)


def bullets(slide, x, y, w, h, items, size=10, color=BLACK, fill=GRAY_PANEL,
            line_color=RGBColor(0x9A, 0x9A, 0x9A), line_w=1.0, dot='•  ',
            leading=1.0, top=0.06):
    """带项目符号的信息面板。items 为字符串列表。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    _fill_line(sp, fill, line_color, line_w)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(0.03)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        run = p.add_run()
        run.text = dot + it
        _set_run_font(run, size, False, color)
    return sp


def line(slide, x1, y1, x2, y2, color=BLACK, w=0.75, arrow='tri',
         dashed=False):
    """带箭头的直线连接符。arrow: 'tri'实心/'none'无。"""
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(x1), Inches(y1),
                                     Inches(x2), Inches(y2))
    cxn.line.color.rgb = color
    cxn.line.width = Pt(w)
    cxn.shadow.inherit = False
    ln = cxn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if arrow == 'tri':
        tail = ln.makeelement(qn('a:tailEnd'),
                              {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return cxn


def save_pdf(prs, pptx_path, pdf_dir):
    import subprocess
    import os
    prs.save(pptx_path)
    subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', pdf_dir, pptx_path],
                   check=True, capture_output=True, timeout=180)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    return os.path.join(pdf_dir, base + '.pdf')
